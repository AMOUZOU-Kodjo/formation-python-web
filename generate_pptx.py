"""Generate a PowerPoint presentation for the 36-module Python training."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

MODULES = [
  {"id": "01-introduction", "title": "Introduction à Python", "short": "Introduction", "part": 1, "duration": "2h", "desc": "Installez Python et écrivez vos premiers scripts.", "icon": "🐍", "color": "#4CAF50"},
  {"id": "02-variables-types", "title": "Variables, types et opérateurs", "short": "Variables & Types", "part": 1, "duration": "2h", "desc": "Maîtrisez les types fondamentaux et les opérateurs.", "icon": "🔢", "color": "#66BB6A"},
  {"id": "03-strings", "title": "Chaînes de caractères", "short": "Strings", "part": 1, "duration": "2h", "desc": "Manipulez le texte : méthodes, slicing, f-strings.", "icon": "📝", "color": "#81C784"},
  {"id": "04-listes-tuples", "title": "Listes et tuples", "short": "Listes & Tuples", "part": 1, "duration": "2h", "desc": "Collections ordonnées : création, méthodes, slicing.", "icon": "📋", "color": "#A5D6A7"},
  {"id": "05-dictionnaires-ensembles", "title": "Dictionnaires et ensembles", "short": "Dicts & Sets", "part": 1, "duration": "2h", "desc": "Associations clé-valeur et ensembles uniques.", "icon": "🗂️", "color": "#C8E6C9"},
  {"id": "06-controle-flux", "title": "Contrôle de flux", "short": "Contrôle de flux", "part": 1, "duration": "2h", "desc": "if/elif/else, match/case, opérateur ternaire.", "icon": "🔀", "color": "#2E7D32"},
  {"id": "07-boucles", "title": "Boucles", "short": "Boucles", "part": 1, "duration": "2h", "desc": "Itérez avec for/while, break/continue/else.", "icon": "🔄", "color": "#388E3C"},
  {"id": "08-fonctions", "title": "Fonctions", "short": "Fonctions", "part": 1, "duration": "3h", "desc": "Fonctions réutilisables, paramètres, portée, closures.", "icon": "⚙️", "color": "#43A047"},
  {"id": "09-comprehensions", "title": "Compréhensions", "short": "Compréhensions", "part": 1, "duration": "2h", "desc": "Code concis avec list/dict/set comprehensions.", "icon": "💡", "color": "#4CAF50"},
  {"id": "10-erreurs-exceptions", "title": "Gestion des erreurs", "short": "Exceptions", "part": 1, "duration": "2h", "desc": "try/except, créez vos propres exceptions.", "icon": "🛡️", "color": "#66BB6A"},
  {"id": "11-modules-packages", "title": "Modules et packages", "short": "Modules", "part": 2, "duration": "2h", "desc": "Organisez votre code en modules et packages.", "icon": "📦", "color": "#FF9800"},
  {"id": "12-fichiers", "title": "Fichiers", "short": "Fichiers", "part": 2, "duration": "2h", "desc": "Lisez/écrivez fichiers, CSV, JSON. pathlib.", "icon": "📁", "color": "#FFA726"},
  {"id": "13-poo-classes", "title": "POO — Classes", "short": "Classes", "part": 2, "duration": "3h", "desc": "Classes, attributs, méthodes, propriétés.", "icon": "🏛️", "color": "#FFB74D"},
  {"id": "14-heritage-polymorphisme", "title": "Héritage et polymorphisme", "short": "Héritage", "part": 2, "duration": "3h", "desc": "Héritage, polymorphisme, classes abstraites.", "icon": "🧬", "color": "#FFCC80"},
  {"id": "15-methodes-speciales", "title": "Méthodes spéciales", "short": "Dunder methods", "part": 2, "duration": "2h", "desc": "Personnalisez avec __str__, __add__, __call__.", "icon": "✨", "color": "#FFE0B2"},
  {"id": "16-iterateurs-generateurs", "title": "Itérateurs et générateurs", "short": "Itérateurs", "part": 2, "duration": "2h", "desc": "Séquences à la demande avec yield.", "icon": "⏳", "color": "#F57C00"},
  {"id": "17-decorateurs", "title": "Décorateurs", "short": "Décorateurs", "part": 2, "duration": "2h", "desc": "Étendez le comportement des fonctions.", "icon": "🎀", "color": "#E65100"},
  {"id": "18-context-managers", "title": "Context managers", "short": "Context managers", "part": 2, "duration": "2h", "desc": "Gérez les ressources avec with/__enter__.", "icon": "📎", "color": "#BF360C"},
  {"id": "19-bibliotheque-standard", "title": "Bibliothèque standard", "short": "Biblio standard", "part": 3, "duration": "3h", "desc": "math, random, datetime, os, sys, json...", "icon": "📚", "color": "#9C27B0"},
  {"id": "20-regex", "title": "Expressions régulières", "short": "Regex", "part": 3, "duration": "3h", "desc": "Regex pour valider, rechercher, transformer.", "icon": "🔍", "color": "#AB47BC"},
  {"id": "21-fonctionnelle", "title": "Programmation fonctionnelle", "short": "Fonctionnelle", "part": 3, "duration": "2h", "desc": "Map, filter, reduce, lambda, itertools.", "icon": "λ", "color": "#BA68C8"},
  {"id": "22-type-hints", "title": "Type hints", "short": "Typage", "part": 3, "duration": "2h", "desc": "Typage avec typing, TypeVar, mypy.", "icon": "🏷️", "color": "#CE93D8"},
  {"id": "23-async-await", "title": "Async / Await", "short": "Async", "part": 3, "duration": "3h", "desc": "Asynchrone : asyncio, gather, aiohttp.", "icon": "⚡", "color": "#E1BEE7"},
  {"id": "24-threading-multiprocess", "title": "Threading & Multiprocess", "short": "Threading", "part": 3, "duration": "3h", "desc": "Parallélisme, GIL, synchronisation.", "icon": "🧵", "color": "#4A148C"},
  {"id": "25-metaclasses-introspection", "title": "Métaclasses & Introspection", "short": "Métaclasses", "part": 4, "duration": "3h", "desc": "Métaclasses, __init_subclass__, introspection.", "icon": "🪞", "color": "#D32F2F"},
  {"id": "26-design-patterns", "title": "Design patterns", "short": "Patterns", "part": 4, "duration": "3h", "desc": "Singleton, Factory, Observer, Strategy.", "icon": "🏗️", "color": "#E53935"},
  {"id": "27-tests-unitaires", "title": "Tests unitaires", "short": "Tests", "part": 4, "duration": "3h", "desc": "pytest, unittest, mocking, TDD.", "icon": "🧪", "color": "#EF5350"},
  {"id": "28-logging-debugging", "title": "Logging & Debugging", "short": "Logging", "part": 4, "duration": "2h", "desc": "Logging, debugging pdb, profiling cProfile.", "icon": "🐛", "color": "#F44336"},
  {"id": "29-environnements-virtuels", "title": "Environnements virtuels", "short": "Venv", "part": 4, "duration": "2h", "desc": "venv, pip, poetry, conda.", "icon": "📦", "color": "#EF9A9A"},
  {"id": "30-packaging-distribution", "title": "Packaging & Distribution", "short": "Packaging", "part": 4, "duration": "2h", "desc": "Distribuez sur PyPI avec pyproject.toml.", "icon": "📤", "color": "#FFCDD2"},
  {"id": "31-api-rest", "title": "API REST avec FastAPI", "short": "API REST", "part": 4, "duration": "4h", "desc": "APIs REST avec FastAPI, Pydantic, Uvicorn.", "icon": "🌐", "color": "#1565C0"},
  {"id": "32-bases-de-donnees", "title": "Bases de données", "short": "BDD", "part": 4, "duration": "4h", "desc": "SQLite, SQLAlchemy ORM, migrations.", "icon": "🗄️", "color": "#1976D2"},
  {"id": "33-web-scraping", "title": "Web scraping", "short": "Scraping", "part": 4, "duration": "3h", "desc": "Requests, BeautifulSoup, Selenium.", "icon": "🕸️", "color": "#1E88E5"},
  {"id": "34-data-science", "title": "Data Science", "short": "Data Science", "part": 4, "duration": "4h", "desc": "NumPy, Pandas : DataFrames, agrégation.", "icon": "📊", "color": "#42A5F5"},
  {"id": "35-visualisation", "title": "Visualisation", "short": "Visualisation", "part": 4, "duration": "3h", "desc": "Graphiques avec Matplotlib et Seaborn.", "icon": "📈", "color": "#64B5F6"},
  {"id": "36-projet-final", "title": "Projet final", "short": "Projet final", "part": 4, "duration": "8h", "desc": "Projet complet pour valider toutes les compétences.", "icon": "🏆", "color": "#0D47A1"},
]

PARTS = [
  {"id": 1, "title": "Les Fondamentaux", "color": "#4CAF50", "desc": "10 modules — Les bases essentielles de Python"},
  {"id": 2, "title": "Intermédiaire", "color": "#FF9800", "desc": "8 modules — Approfondissez vos connaissances"},
  {"id": 3, "title": "Avancé", "color": "#9C27B0", "desc": "6 modules — Concepts avancés du langage"},
  {"id": 4, "title": "Expert / Spécialisation", "color": "#D32F2F", "desc": "12 modules — Maîtrisez l'écosystème Python"},
]

BLUE = RGBColor(0x30, 0x69, 0x98)
GOLD = RGBColor(0xFF, 0xE8, 0x73)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1E, 0x1E, 0x1E)
GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT = RGBColor(0xF5, 0xF5, 0xF5)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def hex_to_rgb(h):
  h = h.lstrip('#')
  return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

def add_bg(slide, color):
  bg = slide.background
  fill = bg.fill
  fill.solid()
  fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, color, alpha=None):
  from pptx.oxml.ns import qn
  shape = slide.shapes.add_shape(1, left, top, width, height)
  shape.fill.solid()
  shape.fill.fore_color.rgb = color
  shape.line.fill.background()
  if alpha is not None:
    shape.fill.fore_color.brightness = 0
  return shape

def add_textbox(slide, left, top, width, height, text, size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT, font_name="Calibri"):
  txBox = slide.shapes.add_textbox(left, top, width, height)
  tf = txBox.text_frame
  tf.word_wrap = True
  p = tf.paragraphs[0]
  p.text = text
  p.font.size = Pt(size)
  p.font.bold = bold
  p.font.color.rgb = color
  p.font.name = font_name
  p.alignment = align
  return txBox

# ---- Slide 1: Title ----
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BLUE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.15), GOLD)
add_shape(slide, Inches(0), Inches(7.35), Inches(13.333), Inches(0.15), GOLD)
add_textbox(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(1), "FORMATION PYTHON", 48, True, WHITE, PP_ALIGN.CENTER)
add_textbox(slide, Inches(1.5), Inches(2.6), Inches(10), Inches(1), "Du débutant à l'expert", 28, False, GOLD, PP_ALIGN.CENTER)
add_textbox(slide, Inches(1.5), Inches(3.6), Inches(10), Inches(0.8), "36 modules — 4 parties — 90+ heures de formation", 20, False, WHITE, PP_ALIGN.CENTER)
add_textbox(slide, Inches(1.5), Inches(5.5), Inches(10), Inches(0.6), "© Formation Python — formation-python-web.vercel.app", 14, False, GOLD, PP_ALIGN.CENTER)

# ---- Slides: Parts Overview ----
for part in PARTS:
  slide = prs.slides.add_slide(prs.slide_layouts[6])
  color = hex_to_rgb(part["color"])
  add_bg(slide, color)
  add_shape(slide, Inches(0), Inches(0), Inches(0.25), Inches(7.5), WHITE)
  add_textbox(slide, Inches(1.5), Inches(1), Inches(10), Inches(1), f"Partie {part['id']}", 22, False, WHITE, PP_ALIGN.LEFT)
  add_textbox(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(1.5), part["title"], 44, True, WHITE, PP_ALIGN.LEFT)
  add_textbox(slide, Inches(1.5), Inches(3.2), Inches(10), Inches(0.8), part["desc"], 20, False, WHITE, PP_ALIGN.LEFT)

# ---- Slides: Modules ----
for mod in MODULES:
  slide = prs.slides.add_slide(prs.slide_layouts[6])
  part_color = hex_to_rgb(PARTS[mod["part"] - 1]["color"])

  # Left color bar
  add_shape(slide, Inches(0), Inches(0), Inches(0.35), Inches(7.5), part_color)

  # Top bar
  add_shape(slide, Inches(0.35), Inches(0), Inches(12.983), Inches(0.08), GOLD)

  # Part tag
  add_textbox(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.6),
              f"Partie {mod['part']} — {PARTS[mod['part'] - 1]['title']}", 14, False, part_color)

  # Module number
  num = mod["id"].split("-")[0]
  add_textbox(slide, Inches(0.8), Inches(1.1), Inches(2), Inches(1),
              f"Module {num}", 16, True, GRAY)

  # Icon
  add_textbox(slide, Inches(10.5), Inches(0.8), Inches(2), Inches(1.2),
              mod["icon"], 48, False, DARK, PP_ALIGN.RIGHT)

  # Title
  add_textbox(slide, Inches(0.8), Inches(2.2), Inches(11), Inches(1.2),
              mod["title"], 36, True, DARK, PP_ALIGN.LEFT)

  # Divider line
  add_shape(slide, Inches(0.8), Inches(3.4), Inches(4), Inches(0.04), part_color)

  # Description
  add_textbox(slide, Inches(0.8), Inches(3.8), Inches(11), Inches(1),
              mod["desc"], 20, False, GRAY, PP_ALIGN.LEFT)

  # Duration badge
  add_textbox(slide, Inches(0.8), Inches(5.2), Inches(11), Inches(0.6),
              f"Durée : {mod['duration']}", 18, False, part_color)

  # Bottom bar
  add_shape(slide, Inches(0.35), Inches(7.42), Inches(12.983), Inches(0.08), GOLD)

# ---- Last Slide: Conclusion ----
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BLUE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.15), GOLD)
add_shape(slide, Inches(0), Inches(7.35), Inches(13.333), Inches(0.15), GOLD)
add_textbox(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(1), "Félicitations !", 48, True, WHITE, PP_ALIGN.CENTER)
add_textbox(slide, Inches(1.5), Inches(2.8), Inches(10), Inches(1), "Vous avez terminé les 36 modules", 28, False, GOLD, PP_ALIGN.CENTER)
add_textbox(slide, Inches(1.5), Inches(4), Inches(10), Inches(0.8), "Obtenez votre attestation sur formation-python-web.vercel.app/certificate", 18, False, WHITE, PP_ALIGN.CENTER)

# Save
prs.save("FormationPython_Complete.pptx")
print("Presentation generated: FormationPython_Complete.pptx")
print(f"  Slides: {len(prs.slides)}")
