"""Generate video presentation PowerPoints for ALL 36 modules with auto-generated content."""

import sys, re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

GOLD = RGBColor(0xFF, 0xE8, 0x73)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1E, 0x1E, 0x1E)
GRAY = RGBColor(0x66, 0x66, 0x66)

MODULES = [
  {"id": "01-introduction", "title": "Introduction à Python", "part": 1, "duration": "2h", "desc": "Installez Python et écrivez vos premiers scripts. Découvrez l'histoire, l'écosystème et les outils essentiels.", "icon": "🐍", "color": "#4CAF50"},
  {"id": "02-variables-types", "title": "Variables, types et opérateurs", "part": 1, "duration": "2h", "desc": "Maîtrisez les types fondamentaux (int, float, str, bool), les variables et tous les opérateurs Python.", "icon": "🔢", "color": "#66BB6A"},
  {"id": "03-strings", "title": "Chaînes de caractères", "part": 1, "duration": "2h", "desc": "Manipulez le texte comme un pro : méthodes, slicing, f-strings, formatage avancé.", "icon": "📝", "color": "#81C784"},
  {"id": "04-listes-tuples", "title": "Listes et tuples", "part": 1, "duration": "2h", "desc": "Les collections ordonnées : création, méthodes, slicing, compréhensions simples.", "icon": "📋", "color": "#A5D6A7"},
  {"id": "05-dictionnaires-ensembles", "title": "Dictionnaires et ensembles", "part": 1, "duration": "2h", "desc": "Associations clé-valeur et collections d'éléments uniques avec leurs cas d'usage.", "icon": "🗂️", "color": "#C8E6C9"},
  {"id": "06-controle-flux", "title": "Contrôle de flux", "part": 1, "duration": "2h", "desc": "Maîtrisez if/elif/else, match/case, l'opérateur ternaire et le walrus operator.", "icon": "🔀", "color": "#2E7D32"},
  {"id": "07-boucles", "title": "Boucles", "part": 1, "duration": "2h", "desc": "Itérez efficacement avec for et while, contrôlez le flux avec break/continue/else.", "icon": "🔄", "color": "#388E3C"},
  {"id": "08-fonctions", "title": "Fonctions", "part": 1, "duration": "3h", "desc": "Créez des fonctions réutilisables, maîtrisez les paramètres, la portée et les closures.", "icon": "⚙️", "color": "#43A047"},
  {"id": "09-comprehensions", "title": "Compréhensions", "part": 1, "duration": "2h", "desc": "Écrivez du code concis et performant avec les list/dict/set comprehensions.", "icon": "💡", "color": "#4CAF50"},
  {"id": "10-erreurs-exceptions", "title": "Gestion des erreurs", "part": 1, "duration": "2h", "desc": "Anticipez et gérez les erreurs avec try/except, créez vos propres exceptions.", "icon": "🛡️", "color": "#66BB6A"},
  {"id": "11-modules-packages", "title": "Modules et packages", "part": 2, "duration": "2h", "desc": "Organisez votre code en modules et packages réutilisables.", "icon": "📦", "color": "#FF9800"},
  {"id": "12-fichiers", "title": "Fichiers", "part": 2, "duration": "2h", "desc": "Lisez et écrivez des fichiers texte, CSV, JSON. Maîtrisez pathlib et les context managers.", "icon": "📁", "color": "#FFA726"},
  {"id": "13-poo-classes", "title": "POO — Classes", "part": 2, "duration": "3h", "desc": "Plongez dans la programmation orientée objet : classes, attributs, méthodes, propriétés.", "icon": "🏛️", "color": "#FFB74D"},
  {"id": "14-heritage-polymorphisme", "title": "Héritage et polymorphisme", "part": 2, "duration": "3h", "desc": "Héritage, polymorphisme, classes abstraites, héritage multiple et MRO.", "icon": "🧬", "color": "#FFCC80"},
  {"id": "15-methodes-speciales", "title": "Méthodes spéciales", "part": 2, "duration": "2h", "desc": "Personnalisez vos objets avec __str__, __add__, __call__ et 50+ autres méthodes.", "icon": "✨", "color": "#FFE0B2"},
  {"id": "16-iterateurs-generateurs", "title": "Itérateurs et générateurs", "part": 2, "duration": "2h", "desc": "Créez des séquences à la demande avec yield et économisez la mémoire.", "icon": "⏳", "color": "#F57C00"},
  {"id": "17-decorateurs", "title": "Décorateurs", "part": 2, "duration": "2h", "desc": "Étendez le comportement des fonctions sans les modifier avec les décorateurs.", "icon": "🎀", "color": "#E65100"},
  {"id": "18-context-managers", "title": "Context managers", "part": 2, "duration": "2h", "desc": "Gérez automatiquement les ressources avec with, __enter__/__exit__ et contextlib.", "icon": "📎", "color": "#BF360C"},
  {"id": "19-bibliotheque-standard", "title": "Bibliothèque standard", "part": 3, "duration": "3h", "desc": "Explorez math, random, datetime, os, sys, json, collections et plus.", "icon": "📚", "color": "#9C27B0"},
  {"id": "20-regex", "title": "Expressions régulières", "part": 3, "duration": "3h", "desc": "Maîtrisez les regex pour valider, rechercher et transformer du texte.", "icon": "🔍", "color": "#AB47BC"},
  {"id": "21-fonctionnelle", "title": "Programmation fonctionnelle", "part": 3, "duration": "2h", "desc": "Map, filter, reduce, lambda, itertools et functools pour un code élégant.", "icon": "λ", "color": "#BA68C8"},
  {"id": "22-type-hints", "title": "Type hints", "part": 3, "duration": "2h", "desc": "Ajoutez des types à votre code avec typing, TypeVar, Protocol et validez avec mypy.", "icon": "🏷️", "color": "#CE93D8"},
  {"id": "23-async-await", "title": "Async / Await", "part": 3, "duration": "3h", "desc": "Programmation asynchrone : asyncio, gather, create_task et aiohttp.", "icon": "⚡", "color": "#E1BEE7"},
  {"id": "24-threading-multiprocess", "title": "Threading & Multiprocess", "part": 3, "duration": "3h", "desc": "Parallélisme avec threading et multiprocessing, GIL, synchronisation.", "icon": "🧵", "color": "#4A148C"},
  {"id": "25-metaclasses-introspection", "title": "Métaclasses & Introspection", "part": 4, "duration": "3h", "desc": "Introspection avancée, métaclasses, __init_subclass__, cas concrets (ORM, DSL).", "icon": "🪞", "color": "#D32F2F"},
  {"id": "26-design-patterns", "title": "Design patterns", "part": 4, "duration": "3h", "desc": "Singleton, Factory, Observer, Strategy et autres patterns en Python.", "icon": "🏗️", "color": "#E53935"},
  {"id": "27-tests-unitaires", "title": "Tests unitaires", "part": 4, "duration": "3h", "desc": "Testez votre code avec pytest, unittest, mocking et TDD.", "icon": "🧪", "color": "#EF5350"},
  {"id": "28-logging-debugging", "title": "Logging & Debugging", "part": 4, "duration": "2h", "desc": "Journalisation avancée, debugging avec pdb, profiling avec cProfile.", "icon": "🐛", "color": "#F44336"},
  {"id": "29-environnements-virtuels", "title": "Environnements virtuels", "part": 4, "duration": "2h", "desc": "venv, pip, pipenv, poetry, conda et bonnes pratiques de gestion de projet.", "icon": "📦", "color": "#EF9A9A"},
  {"id": "30-packaging-distribution", "title": "Packaging & Distribution", "part": 4, "duration": "2h", "desc": "Créez et distribuez vos packages Python sur PyPI avec pyproject.toml.", "icon": "📤", "color": "#FFCDD2"},
  {"id": "31-api-rest", "title": "API REST avec FastAPI", "part": 4, "duration": "4h", "desc": "Créez des APIs REST modernes avec FastAPI, Pydantic et Uvicorn.", "icon": "🌐", "color": "#1565C0"},
  {"id": "32-bases-de-donnees", "title": "Bases de données", "part": 4, "duration": "4h", "desc": "SQLite, SQLAlchemy ORM, relations, sessions, migrations Alembic.", "icon": "🗄️", "color": "#1976D2"},
  {"id": "33-web-scraping", "title": "Web scraping", "part": 4, "duration": "3h", "desc": "Extrayez des données du web avec requests, BeautifulSoup et Selenium.", "icon": "🕸️", "color": "#1E88E5"},
  {"id": "34-data-science", "title": "Data Science", "part": 4, "duration": "4h", "desc": "Analysez des données avec NumPy et Pandas : vectorisation, DataFrames, agrégation.", "icon": "📊", "color": "#42A5F5"},
  {"id": "35-visualisation", "title": "Visualisation", "part": 4, "duration": "3h", "desc": "Créez des graphiques percutants avec Matplotlib et Seaborn.", "icon": "📈", "color": "#64B5F6"},
  {"id": "36-projet-final", "title": "Projet final", "part": 4, "duration": "8h", "desc": "Appliquez toutes vos compétences sur un projet complet au choix.", "icon": "🏆", "color": "#0D47A1"},
]

PARTS = {1: "Les Fondamentaux", 2: "Intermédiaire", 3: "Avancé", 4: "Expert / Spécialisation"}

# ---- Manual slide content for modules that need specific data ----
MANUAL_SLIDES = {
  "01-introduction": [
    ("Qu'est-ce que Python ?", [
      "Crée en 1991 par Guido van Rossum",
      "Langage simple, lisible et polyvalent",
      "Utilisé pour le web, data science, IA, automatisation",
      "Grande communauté, open-source, gratuit"
    ], "print(\"Hello, World!\")\n\n# Python est un langage interprété\n# Pas besoin de compilation\n# Il suffit d'avoir Python installe"),
    ("Installation", [
      "python.org → Download → Installer",
      "Sous Windows : cocher \"Add Python to PATH\"",
      "Verifier avec : python --version",
      "Editeur recommande : VS Code ou PyCharm"
    ], ""),
    ("Premier script", [
      "Creer un fichier .py avec un editeur",
      "Executer avec : python fichier.py",
      "print() pour afficher",
      "input() pour lire une saisie"
    ], "nom = input(\"Quel est ton prenom ? \")\nprint(f\"Bonjour {nom} !\")"),
  ],
  "02-variables-types": [
    ("Les variables", [
      "Une variable stocke une valeur sous un nom",
      "Python deduit le type automatiquement",
      "Snake_case : prenom, age_utilisateur",
      "Sensible a la casse : age != Age"
    ], "prenom = \"Alice\"\nage = 25\ntaille = 1.68\nest_etudiant = True"),
    ("Types fondamentaux", [
      "int : nombres entiers (42)",
      "float : nombres decimaux (3.14)",
      "str : chaines de caracteres (\"Hello\")",
      "bool : booleens (True/False)"
    ], "type(42)      # <class 'int'>\ntype(3.14)    # <class 'float'>\ntype(\"Hello\") # <class 'str'>\ntype(True)    # <class 'bool'>"),
    ("Operateurs", [
      "Arithmetiques : +, -, *, /, //, %, **",
      "Comparaison : ==, !=, >, <, >=, <=",
      "Logiques : and, or, not",
      "Conversion : int(), float(), str(), bool()"
    ], "print(10 / 3)   # 3.333...\nprint(10 // 3)  # 3\nprint(10 % 3)   # 1\nprint(2 ** 8)   # 256"),
  ],
  "03-strings": [
    ("Creation de chaines", [
      "Guillemets simples '...' ou doubles \"...\"",
      "Triples guillemets pour texte multi-lignes",
      "Caracteres d'echappement : \\n, \\t, \\\\"
    ], "s1 = 'Hello'\ns2 = \"World\"\ns3 = \"\"\"Texte\nmulti-lignes\"\"\""),
    ("Indexation et slicing", [
      "Acces par index : texte[0], texte[-1]",
      "Slicing : texte[debut:fin:pas]",
      "texte[::-1] inverse la chaine"
    ], "s = \"Python\"\nprint(s[0])    # P\nprint(s[-1])   # n\nprint(s[0:3])  # Pyt\nprint(s[::-1]) # nohtyP"),
    ("Methodes essentielles", [
      "upper(), lower(), strip()",
      "replace(), split(), join()",
      "startswith(), isdigit(), isalpha()",
      "f-strings pour un affichage moderne"
    ], "s = \"  Hello World  \"\nprint(s.strip())\nprint(s.upper())\nprint(s.split())\nprint(\"-\".join([\"a\",\"b\"]))"),
  ],
  "04-listes-tuples": [
    ("Les listes", [
      "Collection ordonnee et modifiable",
      "Creer : ma_liste = [1, 2, 3]",
      "Indexation et slicing (comme les strings)",
      "Methodes : append(), remove(), sort()"
    ], "fruits = [\"pomme\", \"banane\", \"cerise\"]\nfruits.append(\"datte\")\nfruits.remove(\"banane\")\nprint(fruits[0])  # pomme\nprint(fruits[-1]) # datte"),
    ("Les tuples", [
      "Collection ordonnee et immuable",
      "Creer : mon_tuple = (1, 2, 3)",
      "Plus rapide qu'une liste",
      "Utilise pour des donnees qui ne changent pas"
    ], "coords = (10, 20)\nx, y = coords  # unpacking\nprint(x)  # 10\n\n# Immutable\n# coords[0] = 5  # ERREUR !"),
    ("Methodes courantes", [
      "len() : taille de la collection",
      "min(), max(), sum() pour les nombres",
      "in : verifier la presence d'un element",
      "Indexage negatif et slicing"
    ], "notes = [15, 12, 18, 10, 14]\nprint(len(notes))    # 5\nprint(sum(notes))    # 69\nprint(min(notes))    # 10\nprint(15 in notes)   # True"),
  ],
  "05-dictionnaires-ensembles": [
    ("Les dictionnaires", [
      "Association cle → valeur",
      "Equivalent du JSON en Python",
      "Acces rapide par cle",
      "Modifiable et dynamique"
    ], "eleve = {\n    \"nom\": \"Alice\",\n    \"age\": 25,\n    \"notes\": [15, 18]\n}\nprint(eleve[\"nom\"])  # Alice\neleve[\"age\"] = 26"),
    ("Methodes des dictionnaires", [
      "keys(), values(), items()",
      "get() pour un acces securise",
      "pop(), update(), clear()",
      "Comprehension de dictionnaire"
    ], "for k, v in eleve.items():\n    print(f\"{k}: {v}\")\n\n# Acces securise\nprint(eleve.get(\"prenom\", \"inconnu\"))"),
    ("Les ensembles (set)", [
      "Collection non ordonnee sans doublons",
      "Creer : mon_set = {1, 2, 3}",
      "Operations : union, intersection, difference",
      "Ideal pour enlever les doublons"
    ], "a = {1, 2, 3, 4}\nb = {3, 4, 5, 6}\nprint(a | b)  # union {1,2,3,4,5,6}\nprint(a & b)  # intersection {3,4}\nprint(a - b)  # difference {1,2}"),
  ],
  "06-controle-flux": [
    ("if / elif / else", [
      "Structure conditionnelle de base",
      "Pas de parentheses obligatoires",
      "Indentation definit les blocs",
      "elif = else if"
    ], "age = 18\nif age < 12:\n    print(\"Enfant\")\nelif age < 18:\n    print(\"Adolescent\")\nelse:\n    print(\"Adulte\")"),
    ("Operateurs speciaux", [
      "Operateur ternaire : x if cond else y",
      "match/case (Python 3.10+)",
      "Walrus operator := (assigner + tester)",
      "Combinaison avec and, or, not"
    ], "# Ternaire\nstatut = \"majeur\" if age >= 18 else \"mineur\"\n\n# Walrus\nif (n := len(liste)) > 5:\n    print(f\"Grande liste: {n}\")"),
    ("Bonnes pratiques", [
      "Eviter les if imbriques trop profonds",
      "Utiliser des fonctions pour clarifier",
      "Preferer match/case pour multi-conditions",
      "and/or court-circuitent l'evaluation"
    ], "# Plutot que if imbriques\ndef categorie(age):\n    if age < 12: return \"Enfant\"\n    if age < 18: return \"Ado\"\n    return \"Adulte\""),
  ],
  "07-boucles": [
    ("La boucle for", [
      "Iterer sur une collection",
      "for element in collection:",
      "range() pour une sequence numerique",
      "enumerate() pour index + valeur"
    ], "for i in range(5):\n    print(i)  # 0 1 2 3 4\n\nfruits = [\"pomme\", \"banane\"]\nfor i, f in enumerate(fruits):\n    print(f\"{i}: {f}\")"),
    ("La boucle while", [
      "Repeter tant qu'une condition est vraie",
      "Attention aux boucles infinies !",
      "Utiliser un compteur ou une condition",
      "break pour sortir, continue pour sauter"
    ], "compteur = 0\nwhile compteur < 5:\n    print(compteur)\n    compteur += 1\n\n# Boucle infinie volontaire\nwhile True:\n    reponse = input(\"q pour quitter: \")\n    if reponse == \"q\":\n        break"),
    ("break / continue / else", [
      "break : sort de la boucle",
      "continue : passe a l'iteration suivante",
      "else : s'execute si pas de break",
      "Utilisable avec for et while"
    ], "for n in range(10):\n    if n == 3:\n        continue  # saute 3\n    if n == 7:\n        break     # s'arrete a 7\n    print(n)  # 0 1 2 4 5 6\n\n# else: pas de break\nfor n in range(3):\n    print(n)\nelse:\n    print(\"Fini!\")"),
  ],
  "08-fonctions": [
    ("Definir une fonction", [
      "def nom_fonction(parametres):",
      "Bloc indente = corps de la fonction",
      "return pour renvoyer une valeur",
      "Sans return : la fonction renvoie None"
    ], "def saluer(nom):\n    return f\"Bonjour {nom}!\"\n\nmessage = saluer(\"Alice\")\nprint(message)  # Bonjour Alice!"),
    ("Parametres avances", [
      "Parametres par defaut : def f(x=10)",
      "Parametres nommes : f(x=5, y=3)",
      "*args : nombre variable d'arguments",
      "**kwargs : arguments nommes variables"
    ], "def somme(*args):\n    return sum(args)\n\nprint(somme(1, 2, 3, 4))  # 10\n\ndef afficher(**kwargs):\n    for k, v in kwargs.items():\n        print(f\"{k} = {v}\")"),
    ("Portee des variables", [
      "Variables locales : dans la fonction",
      "Variables globales : en dehors",
      "global pour modifier une globale",
      "Les closures capturent le contexte"
    ], "x = 10  # globale\n\ndef modifier():\n    global x\n    x = 20\n\nmodifier()\nprint(x)  # 20"),
  ],
  "09-comprehensions": [
    ("List comprehensions", [
      "Syntaxe concise pour creer des listes",
      "[expr for elem in iterable]",
      "Avec filtre : [expr for elem in iterable if cond]",
      "Plus rapide qu'une boucle for classique"
    ], "# Classique\ncarres = []\nfor i in range(10):\n    carres.append(i**2)\n\n# Comprehension\ncarres = [i**2 for i in range(10)]"),
    ("Dict et set comprehensions", [
      "{k: v for k, v in iterable}",
      "{expr for elem in iterable}",
      "Comprehension de tuple avec generator",
      "Applicable partout"
    ], "# Dict comprehension\ncarres = {i: i**2 for i in range(5)}\n# {0:0, 1:1, 2:4, 3:9, 4:16}\n\n# Set comprehension\npairs = {i for i in range(10) if i % 2 == 0}\n# {0, 2, 4, 6, 8}"),
    ("Comprehensions imbriquees", [
      "Plusieurs for dans une comprehension",
      "Equivalent a des boucles imbriquees",
      "A utiliser avec moderation (lisibilite)",
      "Peut remplacer filter() et map()"
    ], "# Double boucle\npaires = [(i, j) for i in range(3) for j in range(2)]\n# [(0,0), (0,1), (1,0), (1,1), (2,0), (2,1)]\n\n# Avec condition\nnombres = [i for i in range(20) if i % 2 == 0 if i % 3 == 0]\n# [0, 6, 12, 18]"),
  ],
  "10-erreurs-exceptions": [
    ("Try / Except", [
      "try : bloc a surveiller",
      "except : gerer l'erreur",
      "finally : toujours execute (optionnel)",
      "else : si pas d'erreur (optionnel)"
    ], "try:\n    x = int(input(\"Nombre: \"))\n    print(10 / x)\nexcept ValueError:\n    print(\"Pas un nombre!\")\nexcept ZeroDivisionError:\n    print(\"Division par zero!\")\nfinally:\n    print(\"Termine\")"),
    ("Creer ses exceptions", [
      "class MonErreur(Exception):",
      "Heriter de la classe Exception",
      "Ajouter un message personnalise",
      "Utiliser raise pour declencher"
    ], "class SoldeInsuffisant(Exception):\n    pass\n\ndef retrait(solde, montant):\n    if montant > solde:\n        raise SoldeInsuffisant(\"Solde insuffisant!\")\n    return solde - montant"),
    ("Bonnes pratiques", [
      "Ne pas capturer Exception generiquement",
      "Preciser le type d'exception",
      "Utiliser finally pour les ressources",
      "Context managers (with) pour les fichiers"
    ], "# Mauvais\ntry:\n    ...\nexcept:  # Trop large\n    pass\n\n# Bon\ntry:\n    ...\nexcept (ValueError, TypeError) as e:\n    print(f\"Erreur: {e}\")"),
  ],
  "11-modules-packages": [
    ("Modules", [
      "Un module = un fichier .py",
      "Reutiliser du code avec import",
      "import module, from module import ...",
      "alias : import module as m"
    ], "# mon_module.py\ndef saluer(nom):\n    return f\"Bonjour {nom}\"\n\n# main.py\nimport mon_module\nprint(mon_module.saluer(\"Alice\"))"),
    ("Packages", [
      "Un package = un dossier avec __init__.py",
      "Organiser des modules en arborescence",
      "Import absolu et relatif",
      "__init__.py peut exporter des noms"
    ], "mon_package/\n    __init__.py\n    module_a.py\n    module_b.py\n\nfrom mon_package import module_a\nfrom mon_package.module_b import fonction"),
    ("Bonnes pratiques", [
      "Un module = une responsabilite",
      "Noms courts, explicites, snake_case",
      "Eviter les imports circulaires",
      "if __name__ == \"__main__\": pour les tests"
    ], "# Testable en tant que script\nif __name__ == \"__main__\":\n    # Code de test ici\n    print(saluer(\"Alice\"))\n\n# Execute seulement si lance directement"),
  ],
}

SLIDE_TEMPLATES = {
  "default": [
    ("Introduction", lambda m: [
      f"Module {m['id'][:2]}: {m['title']}",
      m["desc"],
      f"Duree : {m['duration']}",
      f"Partie : {PARTS[m['part']]}",
    ]),
    ("Concepts cles", lambda m: [
      f"Decouvrez les concepts fondamentaux de {m['title'].lower()}",
      "Cas d'usage concrets et exemples pratiques",
      "Bonnes pratiques et pieges a eviter",
      "Exercices pour valider la comprehension",
    ]),
    ("Exemple pratique", lambda m: [
      "Un exemple concret pour illustrer le module",
      "Code commente pas a pas",
      "Testez par vous-meme dans votre environnement",
    ]),
    ("Bonnes pratiques", lambda m: [
      "Suivez les conventions Python (PEP 8)",
      "Documentez votre code",
      "Testez au fur et a mesure",
      "Pratiquez regulierement",
    ]),
  ]
}

def hex_to_rgb(h):
  return RGBColor(*bytes.fromhex(h.lstrip('#')))

def make_presentation(mod):
  prs = Presentation()
  prs.slide_width = Inches(13.333)
  prs.slide_height = Inches(7.5)
  color = hex_to_rgb(mod["color"])

  def add_shape(slide, l, t, w, h, c):
    shape = slide.shapes.add_shape(1, l, t, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = c
    shape.line.fill.background()

  def add_text(slide, l, t, w, h, text, size=22, bold=False, color=DARK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(l, t, w, h)
    box.text_frame.word_wrap = True
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Calibri"
    p.alignment = align

  def add_code(slide, l, t, w, h, code, size=20):
    if not code.strip():
      return
    box = slide.shapes.add_textbox(l, t, w, h)
    box.text_frame.word_wrap = True
    p = box.text_frame.paragraphs[0]
    p.text = code
    p.font.size = Pt(size)
    p.font.name = "Consolas"
    p.font.color.rgb = DARK
    p.alignment = PP_ALIGN.LEFT

  def add_bullets(slide, l, t, w, h, items, size=24, color=DARK):
    box = slide.shapes.add_textbox(l, t, w, h)
    box.text_frame.word_wrap = True
    for i, item in enumerate(items):
      p = box.text_frame.paragraphs[0] if i == 0 else box.text_frame.add_paragraph()
      p.text = f"  {item}"
      p.font.size = Pt(size)
      p.font.color.rgb = color
      p.font.name = "Calibri"
      p.space_after = Pt(10)

  # ---- Title Slide ----
  slide = prs.slides.add_slide(prs.slide_layouts[6])
  add_shape(slide, Inches(0), Inches(3.2), Inches(13.333), Inches(0.1), GOLD)
  add_text(slide, Inches(1), Inches(0.5), Inches(11), Inches(1),
           f"Module {mod['id'][:2]}     {mod['icon']}", 28, False, GRAY, PP_ALIGN.CENTER)
  add_text(slide, Inches(1), Inches(1.5), Inches(11), Inches(2),
           mod["title"], 52, True, DARK, PP_ALIGN.CENTER)
  add_text(slide, Inches(1), Inches(4.5), Inches(11), Inches(1),
           f"Partie {mod['part']} - {PARTS[mod['part']]}", 22, False, GRAY, PP_ALIGN.CENTER)
  add_text(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.8),
           f"Duree : {mod['duration']}  |  {mod['desc'].split('.')[0]}", 18, False, GRAY, PP_ALIGN.CENTER)

  # ---- Content Slides ----
  mod_id = mod["id"]
  if mod_id in MANUAL_SLIDES:
    slides_data = MANUAL_SLIDES[mod_id]
  else:
    slides_data = []
    for title, gen in SLIDE_TEMPLATES["default"]:
      slides_data.append((title, gen(mod), ""))

  for slide_title, bullets, code_str in slides_data:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape(slide, Inches(0), Inches(0), Inches(0.3), Inches(7.5), color)
    add_shape(slide, Inches(0.3), Inches(0), Inches(13.033), Inches(0.08), GOLD)

    add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
             f"Module {mod['id'][:2]} - {mod['title']}", 18, False, color)

    add_shape(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.05), GOLD)
    add_text(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.9),
             slide_title, 40, True, DARK)

    add_bullets(slide, Inches(0.8), Inches(2.4), Inches(5.5), Inches(4.5), bullets)

    if code_str.strip():
      add_text(slide, Inches(7), Inches(2.2), Inches(5.5), Inches(0.5),
               "Exemple de code", 24, True, color)
      add_code(slide, Inches(7), Inches(2.8), Inches(5.5), Inches(3.5), code_str)

    add_shape(slide, Inches(0.3), Inches(7.42), Inches(13.033), Inches(0.08), GOLD)

  # ---- Conclusion Slide ----
  slide = prs.slides.add_slide(prs.slide_layouts[6])
  add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.1), GOLD)
  add_shape(slide, Inches(0), Inches(7.4), Inches(13.333), Inches(0.1), GOLD)
  add_text(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
           "Recapitulatif", 44, True, DARK, PP_ALIGN.CENTER)
  add_bullets(slide, Inches(2), Inches(3), Inches(9), Inches(3),
              [s[0] for s in slides_data], 28)
  add_text(slide, Inches(1), Inches(6.2), Inches(11), Inches(0.7),
           f"Module {mod['id'][:2]} termine !",
           20, False, color, PP_ALIGN.CENTER)

  safe = re.sub(r'[^a-zA-Z0-9_-]', '', mod['title'].replace(' ', '_'))
  fname = f"Video_Module_{mod['id'][:2]}_{safe}.pptx"
  prs.save(fname)
  print(f"  [OK] {fname} ({len(prs.slides)} slides)")

if __name__ == "__main__":
  targets = sys.argv[1:] if len(sys.argv) > 1 else [m["id"] for m in MODULES]
  print(f"Generating {len(targets)} video presentations...")
  for mid in targets:
    mod = next((m for m in MODULES if m["id"] == mid), None)
    if mod:
      make_presentation(mod)
    else:
      print(f"  [??] Unknown: {mid}")
